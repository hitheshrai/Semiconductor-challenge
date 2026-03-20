#!/usr/bin/env python3
"""Add analysis plots, feasibility, and impact slides to ProblemA_Photonator.pptx."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

INTEL_BLUE = RGBColor(0x00, 0x71, 0xC5)
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)

BASE       = os.path.dirname(os.path.abspath(__file__))
PPTX_IN    = os.path.join(BASE, "ProblemA_Photonator.pptx")
PPTX_OUT   = PPTX_IN
PLOTS_DIR  = os.path.join(BASE, "..", "solution", "output")

MARGIN = Inches(0.5)


# ── helpers ───────────────────────────────────────────────────────────────────

def add_blank_slide(prs):
    layout = prs.slide_layouts[34]  # "Blank" — 0 placeholders
    slide = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    return slide


def reorder_slides(prs, new_order):
    """Reorder slides by list of current 0-based indices."""
    sldIdLst = prs.slides._sldIdLst
    sldIds = list(sldIdLst)
    ordered = [sldIds[i] for i in new_order]
    for el in sldIds:
        sldIdLst.remove(el)
    for el in ordered:
        sldIdLst.append(el)


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_header_bar(slide, title, slide_w):
    bar_h = Inches(0.75)
    add_rect(slide, 0, 0, slide_w, bar_h, INTEL_BLUE)
    add_textbox(slide, MARGIN, Inches(0.12), slide_w - 2 * MARGIN, Inches(0.55),
                title, font_size=24, bold=True, color=WHITE)
    return bar_h


def add_image_scaled(slide, path, left, top, max_w, max_h):
    """Add image preserving aspect ratio, centered in its column."""
    if not os.path.exists(path):
        add_textbox(slide, left, top, max_w, Inches(0.4),
                    f"[missing: {os.path.basename(path)}]",
                    font_size=12, italic=True, color=DARK_GREY)
        return
    pic = slide.shapes.add_picture(path, left, top, width=max_w)
    # scale down if height overflows
    if pic.height > max_h:
        ratio = max_h / pic.height
        pic.width = int(pic.width * ratio)
        pic.height = int(max_h)
    # scale down if width overflows after height adjustment
    if pic.width > max_w:
        ratio = max_w / pic.width
        pic.height = int(pic.height * ratio)
        pic.width = int(max_w)
    # center in column
    pic.left = int(left + (max_w - pic.width) / 2)


def add_bullet_slide(slide, slide_w, slide_h, title, items):
    """Add a header bar + bulleted text body. items: list of (indent, text)."""
    bar_h = add_header_bar(slide, title, slide_w)
    content_w = slide_w - 2 * MARGIN
    top = bar_h + Inches(0.25)

    tb = slide.shapes.add_textbox(MARGIN, top, content_w,
                                  slide_h - top - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, (level, text) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        p.space_before = Pt(4) if level == 0 else Pt(0)
        run = p.add_run()
        bullet = "    \u2022 " if level > 0 else "\u2022 "
        run.text = bullet + text
        run.font.size = Pt(16 if level == 0 else 14)
        run.font.color.rgb = DARK_GREY


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation(PPTX_IN)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    content_w = slide_w - 2 * MARGIN

    col_gap = Inches(0.3)
    col_w = int((content_w - col_gap) / 2)
    right_left = int(MARGIN + col_w + col_gap)

    # ── Slide A: Confusion Matrix + ROC ───────────────────────────────────────
    slide_a = add_blank_slide(prs)   # index 9
    bar_h = add_header_bar(slide_a, "Supporting Analysis: Classification Performance", slide_w)
    top = bar_h + Inches(0.1)
    img_h = int(slide_h - bar_h - Inches(0.7))

    add_image_scaled(slide_a,
                     os.path.join(PLOTS_DIR, "plot_cascade_confusion.png"),
                     MARGIN, top, col_w, img_h)
    add_textbox(slide_a, MARGIN, int(slide_h - Inches(0.45)), col_w, Inches(0.35),
                "Cascade Confusion Matrix", font_size=11, italic=True,
                color=DARK_GREY, align=PP_ALIGN.CENTER)

    add_image_scaled(slide_a,
                     os.path.join(PLOTS_DIR, "plot5_roc_curves.png"),
                     right_left, top, col_w, img_h)
    add_textbox(slide_a, right_left, int(slide_h - Inches(0.45)), col_w, Inches(0.35),
                "Per-Class ROC Curves", font_size=11, italic=True,
                color=DARK_GREY, align=PP_ALIGN.CENTER)

    # ── Slide B: t-SNE + Accuracy vs Occurrence ──────────────────────────────
    slide_b = add_blank_slide(prs)   # index 10
    bar_h = add_header_bar(slide_b, "Supporting Analysis: Embedding Quality", slide_w)
    top = bar_h + Inches(0.1)

    add_image_scaled(slide_b,
                     os.path.join(PLOTS_DIR, "plot6_tsne_embeddings.png"),
                     MARGIN, top, col_w, img_h)
    add_textbox(slide_b, MARGIN, int(slide_h - Inches(0.45)), col_w, Inches(0.35),
                "t-SNE Embedding Space", font_size=11, italic=True,
                color=DARK_GREY, align=PP_ALIGN.CENTER)

    add_image_scaled(slide_b,
                     os.path.join(PLOTS_DIR, "plot3_class_accuracy_vs_occurrence.png"),
                     right_left, top, col_w, img_h)
    add_textbox(slide_b, right_left, int(slide_h - Inches(0.45)), col_w, Inches(0.35),
                "Accuracy vs. Training Set Size", font_size=11, italic=True,
                color=DARK_GREY, align=PP_ALIGN.CENTER)

    # ── Slide C: Feasibility & Deployment ─────────────────────────────────────
    slide_c = add_blank_slide(prs)   # index 11
    add_bullet_slide(slide_c, slide_w, slide_h,
                     "Feasibility & Deployment", [
        (0, "Production-ready: model is 84 MB total — runs on any CUDA GPU or CPU"),
        (0, "Inference: 46 ms/image on GPU (well under the 1-second requirement)"),
        (1, "Cold start with model loading: ~600 ms (one-time cost)"),
        (0, "Scalability: new defect types via few-shot prototype registration — no retraining pipeline needed"),
        (1, "Register a new class from a single labelled image in seconds"),
        (0, "Domain generalization: frozen DINOv2 backbone (pretrained on 142M images) transfers across fabs"),
        (1, "Only class prototypes need updating — not model weights"),
        (0, "Minimal infrastructure: single Python script, no cloud dependency, no MLOps pipeline required"),
    ])

    # ── Slide D: Business Impact ──────────────────────────────────────────────
    slide_d = add_blank_slide(prs)   # index 12
    add_bullet_slide(slide_d, slide_w, slide_h,
                     "Business Impact", [
        (0, "91% defect recall: ~9 in 10 defective chips caught before shipping"),
        (1, "Reduces costly field returns and customer-facing failures"),
        (0, "Asymmetric cost advantage: false alarms are cheap; missed defects are not"),
        (1, "Model is tuned to err on the side of flagging — not passing"),
        (0, "Automation: reduces manual inspection burden, frees expert time for edge cases"),
        (0, "Rapid adaptation: new defect types registered in seconds, not days of retraining"),
        (1, "Production lines can respond to novel failure modes immediately"),
        (0, "Consistent quality: eliminates inspector fatigue and subjective variability"),
    ])

    # ── Reorder ───────────────────────────────────────────────────────────────
    # Current: 0-8 original, 9=A, 10=B, 11=C, 12=D
    # Target:  0-5 (orig), 9, 10 (plots), 6, 7 (orig 7-8), 11, 12 (new), 8 (conclusion)
    new_order = [0, 1, 2, 3, 4, 5, 9, 10, 6, 7, 11, 12, 8]
    reorder_slides(prs, new_order)

    prs.save(PPTX_OUT)
    n = len(prs.slides)
    kb = os.path.getsize(PPTX_OUT) // 1024
    print(f"Saved: {PPTX_OUT}")
    print(f"  {n} slides, {kb} KB")
    print(f"  Slide order:")
    for i, name in enumerate([
        "Title", "Project Description", "Our Solution",
        "Model Architecture", "Handling Class Imbalance", "Results",
        "Analysis: Classification Performance (NEW)",
        "Analysis: Embedding Quality (NEW)",
        "Few-Shot Learning", "Assumptions & Future",
        "Feasibility & Deployment (NEW)",
        "Business Impact (NEW)",
        "Conclusion",
    ], 1):
        print(f"    {i:2d}. {name}")


if __name__ == "__main__":
    main()
