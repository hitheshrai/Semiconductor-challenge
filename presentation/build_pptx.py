"""Build presentation PPTX from slide content using python-pptx."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy, os, re

INTEL_BLUE = RGBColor(0x00, 0x71, 0xC5)
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF0, 0xF0, 0xF0)
BLACK      = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.6)
CONTENT_W = SLIDE_W - 2 * MARGIN

OUTPUT_DIR = os.path.dirname(__file__)
PLOTS_DIR  = os.path.join(OUTPUT_DIR, "..", "solution", "output")

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank


# ── helpers ────────────────────────────────────────────────────────────────

def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    # white background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height,
                text, font_size=18, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_header_bar(slide, title_text, subtitle_text=None):
    """Blue bar at the top with title."""
    bar_h = Inches(1.05)
    add_rect(slide, 0, 0, SLIDE_W, bar_h, INTEL_BLUE)
    add_textbox(slide, MARGIN, Inches(0.1), CONTENT_W, Inches(0.7),
                title_text, font_size=28, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT)
    if subtitle_text:
        add_textbox(slide, MARGIN, Inches(0.72), CONTENT_W, Inches(0.28),
                    subtitle_text, font_size=14, color=RGBColor(0xCC, 0xE5, 0xFF),
                    align=PP_ALIGN.LEFT)
    return bar_h


def body_top(bar_h):
    return bar_h + Inches(0.15)


def add_bullet_block(slide, top, items, font_size=17, indent_size=18,
                     line_spacing_pt=6):
    """Add a list of (level, text) bullet tuples as a single text box."""
    left   = MARGIN
    width  = CONTENT_W
    height = SLIDE_H - top - Inches(0.2)
    txBox  = slide.shapes.add_textbox(left, top, width, height)
    tf     = txBox.text_frame
    tf.word_wrap = True

    for i, (level, text) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = level
        p.space_after = Pt(line_spacing_pt)
        run = p.add_run()
        run.text = ("  " * level + "• " if level > 0 else "• ") + text
        run.font.size = Pt(font_size - level * 2)
        run.font.color.rgb = DARK_GREY


def add_two_col_table(slide, top, headers, rows, col_widths=None, font_size=15):
    """Simple table using text boxes in a grid."""
    ncols = len(headers)
    if col_widths is None:
        col_widths = [CONTENT_W / ncols] * ncols

    row_h = Inches(0.37)
    left  = MARGIN

    # header row
    x = left
    for i, h in enumerate(headers):
        add_rect(slide, x, top, col_widths[i], row_h, INTEL_BLUE)
        add_textbox(slide, x + Inches(0.05), top + Inches(0.05),
                    col_widths[i] - Inches(0.1), row_h - Inches(0.05),
                    h, font_size=font_size, bold=True, color=WHITE)
        x += col_widths[i]

    # data rows
    for ri, row in enumerate(rows):
        x = left
        bg = LIGHT_GREY if ri % 2 == 0 else WHITE
        for ci, cell in enumerate(row):
            add_rect(slide, x, top + row_h * (ri + 1),
                     col_widths[ci], row_h, bg)
            add_textbox(slide,
                        x + Inches(0.05),
                        top + row_h * (ri + 1) + Inches(0.05),
                        col_widths[ci] - Inches(0.1),
                        row_h - Inches(0.05),
                        cell, font_size=font_size - 1, color=DARK_GREY)
            x += col_widths[ci]

    return top + row_h * (len(rows) + 1)


def add_code_block(slide, top, code_text, font_size=12):
    left  = MARGIN
    width = CONTENT_W
    lines = code_text.strip().split("\n")
    height = Pt(font_size + 4) * len(lines) + Inches(0.2)
    height = min(height, SLIDE_H - top - Inches(0.1))
    add_rect(slide, left, top, width, height, LIGHT_GREY)
    txBox = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(0.1),
        width - Inches(0.2), height - Inches(0.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = "Courier New"
        run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    return top + height + Inches(0.1)


def add_image_if_exists(slide, path, top, max_height=Inches(4.5)):
    if not os.path.exists(path):
        return top
    left  = MARGIN + Inches(0.5)
    width = CONTENT_W - Inches(1)
    try:
        pic = slide.shapes.add_picture(path, left, top, width=width)
        # scale to fit max_height
        if pic.height > max_height:
            ratio = max_height / pic.height
            pic.height = int(max_height)
            pic.width  = int(pic.width * ratio)
            pic.left   = int(SLIDE_W / 2 - pic.width / 2)
        return top + pic.height + Inches(0.1)
    except Exception as e:
        add_textbox(slide, MARGIN, top, CONTENT_W, Inches(0.4),
                    f"[image: {os.path.basename(path)}]",
                    font_size=13, italic=True, color=DARK_GREY)
        return top + Inches(0.5)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE DEFINITIONS
# ═══════════════════════════════════════════════════════════════════════════

# ── 1. Title slide ──────────────────────────────────────────────────────────
slide = add_slide()
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, INTEL_BLUE)
add_textbox(slide, MARGIN, Inches(1.8), CONTENT_W, Inches(1.2),
            "Few-Shot Semiconductor Defect Classification",
            font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, MARGIN, Inches(3.1), CONTENT_W, Inches(0.5),
            "ASU / Intel Semiconductor Solutions Challenge 2026 — Problem A",
            font_size=18, color=RGBColor(0xCC, 0xE5, 0xFF), align=PP_ALIGN.CENTER)
add_textbox(slide, MARGIN, Inches(3.7), CONTENT_W, Inches(0.4),
            "Two-Stage Cascade  ·  DINOv2  ·  Prototype Inference",
            font_size=15, italic=True, color=RGBColor(0xAA, 0xCC, 0xFF),
            align=PP_ALIGN.CENTER)
add_textbox(slide, MARGIN, Inches(5.5), CONTENT_W, Inches(0.4),
            "Hithesh Vurukonda  ·  March 2026",
            font_size=14, color=RGBColor(0xCC, 0xE5, 0xFF), align=PP_ALIGN.CENTER)


# ── 2. The Problem ──────────────────────────────────────────────────────────
slide = add_slide()
bar_h = add_header_bar(slide, "The Problem",
                        "Wafer inspection at Intel scale — where misses cost millions")
top = body_top(bar_h)

bullets = [
    (0, "3,778 grayscale wafer images across 9 classes"),
    (0, "Severe class imbalance: 94.5% \"good\",  8–50 samples per defect type"),
    (0, "Must detect all 8 defect types with ≤1 labelled example possible"),
    (0, "Shipping a defective chip as \"good\" is far more costly than a false alarm"),
    (0, "Critical metric: defect recall — not overall accuracy"),
]
add_bullet_block(slide, top, bullets, font_size=19)

col_w = [Inches(2.0), Inches(1.5), Inches(2.0), Inches(1.5), Inches(2.0), Inches(1.5)]
top2 = SLIDE_H - Inches(2.6)
add_two_col_table(
    slide, top2,
    ["Class", "Images", "Class", "Images", "Class", "Images"],
    [
        ["good",     "3,572", "defect5",  "25", "defect4",  "14"],
        ["defect10", "38",    "defect1",  "20", "defect3",  "9"],
        ["defect8",  "42",    "defect2",  "~45","defect9",  "8"],
    ],
    col_widths=col_w, font_size=14
)


# ── 3. Why Standard Classifiers Fail ────────────────────────────────────────
slide = add_slide()
bar_h = add_header_bar(slide, "Why Standard Classifiers Fail")
top = body_top(bar_h)

add_textbox(slide, MARGIN, top, CONTENT_W, Inches(0.35),
            "A single model faces irreconcilable objectives:",
            font_size=18, color=DARK_GREY)
top += Inches(0.4)

code = """\
Objective 1 (Stage 1): "good vs. defective"
  → 94.5% good → model predicts "good" for everything
  → defect recall ≈ 0%

Objective 2 (Stage 2): "which defect type?"
  → 8 rare classes, 8–50 samples each
  → overwhelmed by 3,572 "good" samples"""
top = add_code_block(slide, top, code, font_size=14)

add_textbox(slide, MARGIN, top + Inches(0.1), CONTENT_W, Inches(0.5),
            "Baseline single model result:",
            font_size=18, bold=True, color=INTEL_BLUE)
top += Inches(0.6)

add_two_col_table(
    slide, top,
    ["Metric", "Value"],
    [
        ["Overall Accuracy",   "91.1%"],
        ["Balanced Accuracy",  "0.28"],
        ["Avg. Defect Recall", "~20%"],
    ],
    col_widths=[Inches(4.5), Inches(3)], font_size=16
)

add_textbox(slide, MARGIN, SLIDE_H - Inches(0.55), CONTENT_W, Inches(0.4),
            "High overall accuracy masks near-complete failure to detect defects.",
            font_size=15, italic=True, color=RGBColor(0xCC, 0x00, 0x00))


# ── 4. Two-Stage Cascade Design ─────────────────────────────────────────────
slide = add_slide()
bar_h = add_header_bar(slide, "Solution: Two-Stage Cascade")
top = body_top(bar_h)

add_textbox(slide, MARGIN, top, CONTENT_W, Inches(0.35),
            "Decompose into two tractable, independently optimisable problems.",
            font_size=18, color=DARK_GREY)
top += Inches(0.4)

code = """\
Input Image
    │
    ▼
┌──────────────────────────────────┐
│ Stage 1: Binary Classifier        │
│ "Is this wafer defective?"        │
│  Focal Loss + Balanced Sampler    │
└──────────────┬────────────────────┘
               │
      ┌────────┴────────┐
  defect_prob < 0.65   defect_prob ≥ 0.65
      │                      │
      ▼                      ▼
  predict "good"    ┌──────────────────────────────┐
                    │ Stage 2: Defect-Type Classifier │
                    │ "Which of the 8 defect types?"  │
                    │  Prototype cosine similarity    │
                    └─────────────────────────────────┘"""
add_code_block(slide, top, code, font_size=12)

add_textbox(slide, MARGIN, SLIDE_H - Inches(0.55), CONTENT_W, Inches(0.4),
            "Each stage trains on a balanced dataset with no conflicting objectives.",
            font_size=15, italic=True, color=INTEL_BLUE)


# ── 5. Backbone: DINOv2 ─────────────────────────────────────────────────────
slide = add_slide()
bar_h = add_header_bar(slide, "Backbone: DINOv2 ViT-Small/14")
top = body_top(bar_h)

add_textbox(slide, MARGIN, top, CONTENT_W, Inches(0.35),
            "Pretrained on 142M images using self-supervised DINO + iBOT objectives (Meta AI, 2023)",
            font_size=17, color=DARK_GREY)
top += Inches(0.45)

top = add_two_col_table(
    slide, top,
    ["Property", "Supervised ImageNet", "DINOv2"],
    [
        ["Pretraining data",          "~14M labelled",   "142M unlabelled"],
        ["CLS token quality",         "Moderate",        "Excellent for clustering"],
        ["Few-shot defect benchmarks","Baseline",        "+5–8% bal. acc."],
        ["Cosine similarity alignment","Partial",         "Direct"],
        ["NVIDIA semiconductor AUC",  "—",               "98.5% die-level accuracy"],
    ],
    col_widths=[Inches(3.8), Inches(3.2), Inches(3.2)], font_size=15
)

top += Inches(0.2)
bullets = [
    (0, "CLS token embedding designed to cluster visually similar patches"),
    (0, "Self-supervised objectives align naturally with prototype-based inference"),
    (0, "No labelled data needed for pretraining — scales to 142M images"),
]
add_bullet_block(slide, top, bullets, font_size=17)


# ── 6. Model Architecture ───────────────────────────────────────────────────
slide = add_slide()
bar_h = add_header_bar(slide, "Model Architecture")
top = body_top(bar_h)

code = """\
Input (224×224 grayscale → 3-channel RGB replica)
        │
        ▼
DINOv2 ViT-Small/14  (patch size 14, 384-dim CLS token)
        │
        ▼
Embedding Head:
  Linear(384 → 256) → BatchNorm → ReLU → Dropout(0.35)
  → Linear(256 → 256) → BatchNorm → L2-Normalize
        │
        ▼
   Unit Hypersphere  (‖e‖₂ = 1)
        │
        ▼
Cosine similarity to class prototypes  →  prediction"""
top = add_code_block(slide, top, code, font_size=14)

top += Inches(0.1)
bullets = [
    (0, "L2-normalisation: all embeddings projected onto unit sphere"),
    (0, "Cosine similarity reduces to a dot product — efficiently computable"),
    (0, "Dropout(0.35): critical regularisation with only 8–50 samples per defect class"),
    (0, "Grayscale → RGB: single channel replicated; no information loss"),
]
add_bullet_block(slide, top, bullets, font_size=16)


# ── 7. Training: Stage 1 ────────────────────────────────────────────────────
slide = add_slide()
bar_h = add_header_bar(slide, "Training: Stage 1 — Binary Classifier",
                        "Goal: maximise defect recall while good-chip recall ≥ 80%")
top = body_top(bar_h)

top = add_two_col_table(
    slide, top,
    ["Hyperparameter", "Value", "Rationale"],
    [
        ["Epochs",     "30",              "Sufficient given pretrained backbone"],
        ["Optimizer",  "AdamW",           "Weight decay regularisation"],
        ["Learning rate","1.5e-4",        "Conservative for pretrained backbone"],
        ["Weight decay","0.05",           "Strong regularisation (ViT style)"],
        ["Scheduler",  "Cosine + 10-ep warmup","Smooth convergence"],
        ["Loss",       "Focal Loss (γ=2.0)","Down-weights easy good samples"],
        ["Sampler",    "Balanced (good/defective)","Prevents gradient starvation"],
        ["Threshold τ","0.65",            "Tuned on val; maximises defect recall"],
    ],
    col_widths=[Inches(2.5), Inches(2.5), Inches(5.3)], font_size=14
)

top += Inches(0.2)
add_textbox(slide, MARGIN, top, CONTENT_W, Inches(0.6),
            "Focal Loss: FL(p) = −(1−p)^γ log(p)   A sample predicted at 90% confidence "
            "contributes only (0.1)² = 0.01× the weight of standard cross-entropy — "
            "forcing focus on ambiguous defect examples.",
            font_size=14, italic=True, color=DARK_GREY)


# ── 8. Training: Stage 2 ────────────────────────────────────────────────────
slide = add_slide()
bar_h = add_header_bar(slide, "Training: Stage 2 — Defect-Type Classifier",
                        "Goal: discriminate 8 defect types with no \"good\" class competition")
top = body_top(bar_h)

top = add_two_col_table(
    slide, top,
    ["Hyperparameter", "Value", "Rationale"],
    [
        ["Epochs",       "60",               "Longer — 8-way discrimination"],
        ["Optimizer",    "AdamW",             ""],
        ["Learning rate","1e-4",             "Slightly lower; fine-grained task"],
        ["Weight decay", "0.01",             "Less aggressive; small dataset"],
        ["Scheduler",    "Cosine annealing", ""],
        ["Loss",         "CE + label smooth (ε=0.1)","Prevents overconfidence on tiny classes"],
        ["Sampler",      "Balanced (8 defect classes)","Equal class frequency"],
        ["Inference",    "Prototype cosine similarity","No retraining for new classes"],
    ],
    col_widths=[Inches(2.5), Inches(2.7), Inches(5.1)], font_size=14
)

top += Inches(0.2)
add_textbox(slide, MARGIN, top, CONTENT_W, Inches(0.5),
            "Label smoothing ε=0.1: replaces hard one-hot targets with soft distributions, "
            "regularising the output and producing better-calibrated probabilities when training on 8–50 samples/class.",
            font_size=14, italic=True, color=DARK_GREY)


# ── 9. Few-Shot Inference via Prototypes ────────────────────────────────────
slide = add_slide()
bar_h = add_header_bar(slide, "Few-Shot Inference via TTA-Averaged Prototypes")
top = body_top(bar_h)

code = """\
# 4 augmented views per training image:
# identity,  h-flip,  v-flip,  h+v-flip

for each class c:
    prototype[c] = mean(L2_normalize(embed(aug_i(x)))
                        for all x in class c, all aug_i)

# Inference (single forward pass):
embedding    = L2_normalize(model.get_embedding(image))
prediction   = argmax(cosine_similarity(embedding, prototypes))"""
top = add_code_block(slide, top, code, font_size=14)

top += Inches(0.15)
bullets = [
    (0, "TTA on prototypes: flip-invariant class centroids — stable, representative embeddings"),
    (0, "Single-image inference against stable prototypes outperforms 4-image test TTA"),
    (0, "0.909 balanced accuracy (TTA-protos)  vs.  0.880 (test-time TTA) — test TTA adds noise"),
    (0, "Faster too: 1 forward pass at test time instead of 4"),
]
add_bullet_block(slide, top, bullets, font_size=17)


# ── 10. Few-Shot Extensibility ──────────────────────────────────────────────
slide = add_slide()
bar_h = add_header_bar(slide, "Few-Shot Extensibility — No Retraining Required")
top = body_top(bar_h)

add_textbox(slide, MARGIN, top, CONTENT_W, Inches(0.35),
            "New defect types can be registered at runtime from ≥1 labelled example:",
            font_size=18, color=DARK_GREY)
top += Inches(0.4)

code = """\
# Register a new defect class at runtime
python classify.py image.png --cascade --register new_defect examples/*.png

# Internally:
new_proto = mean(L2_normalize(model.get_embedding(x)) for x in examples)
prototypes["new_defect"] = new_proto   # no model weight update"""
top = add_code_block(slide, top, code, font_size=14)

top += Inches(0.2)
add_textbox(slide, MARGIN, top, CONTENT_W, Inches(0.35),
            "N-shot performance on held-out validation (prototype inference):",
            font_size=17, bold=True, color=INTEL_BLUE)
top += Inches(0.4)

add_two_col_table(
    slide, top,
    ["N-shot", "Accuracy"],
    [["1-shot","80.1%"],["2-shot","80.7%"],["5-shot","80.7%"],
     ["8-shot","80.1%"],["20-shot","79.5%"]],
    col_widths=[Inches(3), Inches(3)], font_size=16
)


# ── 11. Results: Progression ────────────────────────────────────────────────
slide = add_slide()
bar_h = add_header_bar(slide, "Results: Progression of Approaches",
                        "All evaluated on the same 20% stratified validation split (seed=42)")
top = body_top(bar_h)

add_two_col_table(
    slide, top,
    ["Approach", "Overall Acc", "Bal. Acc", "Defect Recall"],
    [
        ["Single EfficientNet-B0 baseline", "91.1%", "0.28",  "~20%"],
        ["+ Tau-norm + logit adjustment",   "85.7%", "0.56",  "52.5%"],
        ["EfficientNet cascade (τ=0.35)",   "85.1%", "0.781", "70.7%"],
        ["ViT + MAE pretraining cascade",   "84.7%", "0.780", "78.0%"],
        ["EfficientNet cascade + test TTA", "87.4%", "0.867", "~87%"],
        ["DINOv2 cascade (τ=0.65)",         "87.4%", "0.881", "87.5%"],
        ["DINOv2 + TTA-protos  ✅ FINAL",   "87.7%", "0.909", "91.25%"],
    ],
    col_widths=[Inches(5.5), Inches(1.9), Inches(1.9), Inches(1.9)], font_size=15
)

top = SLIDE_H - Inches(0.7)
bullets = [
    (0, "Each step addressed a specific, identified root cause — not hyperparameter tweaking"),
]
add_bullet_block(slide, top, bullets, font_size=16)


# ── 12. Final Results ───────────────────────────────────────────────────────
slide = add_slide()
bar_h = add_header_bar(slide, "Final Results",
                        "DINOv2 Cascade · TTA-Averaged Prototypes · τ = 0.65")
top = body_top(bar_h)

# summary numbers
for i, (label, value, note) in enumerate([
    ("Overall Accuracy",   "87.7%",   ""),
    ("Balanced Accuracy",  "0.909",   ""),
    ("Avg. Defect Recall", "91.25%",  "business-critical metric"),
    ("Inference time",     "~46 ms",  "warm GPU, single image"),
]):
    y = top + Inches(0.55) * i
    add_rect(slide, MARGIN, y, Inches(3.5), Inches(0.45), LIGHT_GREY)
    add_textbox(slide, MARGIN + Inches(0.1), y + Inches(0.04),
                Inches(3.3), Inches(0.37), label,
                font_size=15, color=DARK_GREY)
    add_textbox(slide, MARGIN + Inches(3.6), y + Inches(0.04),
                Inches(2.5), Inches(0.37), value,
                font_size=18, bold=True, color=INTEL_BLUE)
    if note:
        add_textbox(slide, MARGIN + Inches(6.2), y + Inches(0.09),
                    Inches(4.5), Inches(0.37), f"({note})",
                    font_size=13, italic=True, color=DARK_GREY)

top += Inches(0.55 * 4 + 0.2)

add_textbox(slide, MARGIN, top, CONTENT_W, Inches(0.35),
            "Per-class recall:", font_size=17, bold=True, color=INTEL_BLUE)
top += Inches(0.38)

add_two_col_table(
    slide, top,
    ["Class", "Train", "Val", "Recall", "Notes"],
    [
        ["defect1",  "16", "4",   "100%", ""],
        ["defect2",  "~36","~9",  "100%", ""],
        ["defect3",  "7",  "2",   "100%", ""],
        ["defect4",  "11", "3",   "100%", ""],
        ["defect5",  "20", "5",   "80%",  ""],
        ["defect8",  "34", "8",   "50%",  "Stage 1 bottleneck — visual overlap with good"],
        ["defect9",  "6",  "1",   "100%", "1 val sample"],
        ["defect10", "30", "8",   "100%", ""],
    ],
    col_widths=[Inches(1.5), Inches(1.0), Inches(0.8), Inches(1.2), Inches(5.8)],
    font_size=13
)


# ── 13. Confusion Matrix ─────────────────────────────────────────────────────
slide = add_slide()
bar_h = add_header_bar(slide, "Cascade Confusion Matrix")
top = body_top(bar_h) + Inches(0.05)
add_image_if_exists(slide,
    os.path.join(PLOTS_DIR, "plot_cascade_confusion.png"),
    top, max_height=Inches(5.5))
add_textbox(slide, MARGIN, SLIDE_H - Inches(0.5), CONTENT_W, Inches(0.38),
            "Strong diagonal — defect types well-separated. Main off-diagonal: defect8 classified as \"good\" by Stage 1.",
            font_size=14, italic=True, color=DARK_GREY, align=PP_ALIGN.CENTER)


# ── 14. t-SNE Embedding Space ───────────────────────────────────────────────
slide = add_slide()
bar_h = add_header_bar(slide, "t-SNE Embedding Space (256-d → 2-d)")
top = body_top(bar_h) + Inches(0.05)
add_image_if_exists(slide,
    os.path.join(PLOTS_DIR, "plot6_tsne_embeddings.png"),
    top, max_height=Inches(5.5))
add_textbox(slide, MARGIN, SLIDE_H - Inches(0.5), CONTENT_W, Inches(0.38),
            "DINOv2 embeddings form tight, separable clusters per defect type — even with 8–50 training samples.",
            font_size=14, italic=True, color=DARK_GREY, align=PP_ALIGN.CENTER)


# ── 15. Class Accuracy vs Occurrence ────────────────────────────────────────
slide = add_slide()
bar_h = add_header_bar(slide, "Classification Accuracy vs. Training Set Size")
top = body_top(bar_h) + Inches(0.05)
add_image_if_exists(slide,
    os.path.join(PLOTS_DIR, "plot3_class_accuracy_vs_occurrence.png"),
    top, max_height=Inches(5.5))
add_textbox(slide, MARGIN, SLIDE_H - Inches(0.5), CONTENT_W, Inches(0.38),
            "Most defect classes achieve 100% recall regardless of training set size. Defect8 is the exception.",
            font_size=14, italic=True, color=DARK_GREY, align=PP_ALIGN.CENTER)


# ── 16. ROC Curves ──────────────────────────────────────────────────────────
slide = add_slide()
bar_h = add_header_bar(slide, "ROC Curves — Per-Class AUC")
top = body_top(bar_h) + Inches(0.05)
add_image_if_exists(slide,
    os.path.join(PLOTS_DIR, "plot5_roc_curves.png"),
    top, max_height=Inches(5.5))
add_textbox(slide, MARGIN, SLIDE_H - Inches(0.5), CONTENT_W, Inches(0.38),
            "High AUC across all defect classes. Defect8 shows the flattest ROC — Stage 1 discrimination difficulty.",
            font_size=14, italic=True, color=DARK_GREY, align=PP_ALIGN.CENTER)


# ── 17. Few-Shot Learning Curve ─────────────────────────────────────────────
slide = add_slide()
bar_h = add_header_bar(slide, "Few-Shot Learning Curve")
top = body_top(bar_h) + Inches(0.05)
add_image_if_exists(slide,
    os.path.join(PLOTS_DIR, "plot4_few_shot_learning_curve.png"),
    top, max_height=Inches(5.5))
add_textbox(slide, MARGIN, SLIDE_H - Inches(0.5), CONTENT_W, Inches(0.38),
            "Performance stable from 1-shot onward — DINOv2 embeddings are robust at any support size.",
            font_size=14, italic=True, color=DARK_GREY, align=PP_ALIGN.CENTER)


# ── 18. Hardware & Timing ───────────────────────────────────────────────────
slide = add_slide()
bar_h = add_header_bar(slide, "Hardware & Timing")
top = body_top(bar_h)

top = add_two_col_table(
    slide, top,
    ["Component", "Specification"],
    [
        ["Training machine",          "DGX Spark, NVIDIA GB10 (128 GB VRAM)"],
        ["MAE pretraining",           "~70 min (300 epochs)"],
        ["Cascade training",          "~85 min (Stage 1: 30 ep, Stage 2: 60 ep)"],
        ["Inference — warm GPU",      "~46 ms per image"],
        ["Inference — cold start",    "~600 ms (model loading)"],
        ["Model checkpoint",          "84 MB (DINOv2 ViT-S/14)"],
        ["Competition requirement",   "≤1 second per image  ✅"],
    ],
    col_widths=[Inches(4.5), Inches(8.0)], font_size=16
)

top += Inches(0.25)
add_textbox(slide, MARGIN, top, CONTENT_W, Inches(0.4),
            "GB10 has compute capability 12.1 (exceeds PyTorch's official max of 12.0) — runs without issue.",
            font_size=14, italic=True, color=DARK_GREY)


# ── 19. Design Decision Summary ─────────────────────────────────────────────
slide = add_slide()
bar_h = add_header_bar(slide, "Design Decision Summary")
top = body_top(bar_h)

add_two_col_table(
    slide, top,
    ["Decision", "Choice", "Why"],
    [
        ["Architecture",    "Two-stage cascade",         "Irreconcilable objectives in single model"],
        ["Backbone",        "DINOv2 ViT-S/14",           "+5–8% bal. acc.; cosine-aligned embeddings"],
        ["Loss (Stage 1)",  "Focal Loss (γ=2.0)",        "Auto-focuses on hard defect examples"],
        ["Loss (Stage 2)",  "CE + label smoothing",      "Prevents overconfidence on small classes"],
        ["Inference",       "Prototype cosine sim.",     "No retraining for new defect types"],
        ["TTA strategy",    "On prototypes, not test",   "Stabilises centroids; avoids test noise"],
        ["Threshold τ",     "0.65",                      "Tuned: max defect recall, good recall ≥ 80%"],
    ],
    col_widths=[Inches(2.5), Inches(2.8), Inches(7.0)], font_size=14
)


# ── 20. Future Work ─────────────────────────────────────────────────────────
slide = add_slide()
bar_h = add_header_bar(slide, "Limitations & Future Work")
top = body_top(bar_h)

add_textbox(slide, MARGIN, top, CONTENT_W, Inches(0.35),
            "Current bottleneck: defect8 (50% recall) — Stage 1 visual overlap with \"good\" chips",
            font_size=17, bold=True, color=RGBColor(0xCC, 0x44, 0x00))
top += Inches(0.42)

add_two_col_table(
    slide, top,
    ["Limitation / Opportunity", "Current", "Proposed Fix"],
    [
        ["defect8 Stage 1 discrimination","50% recall","Per-class τ or dedicated anomaly detector"],
        ["Backbone size",     "ViT-Small/14 (22M params)","ViT-Base (86M params), ~90ms inference"],
        ["Image resolution",  "224×224 (from ≤1500×2500)","Tiled inference at native resolution"],
        ["Static prototypes", "Computed at train time",   "Online EMA updates from production data"],
        ["Val set size",      "1–2 samples/rare class",   "More labelled data improves estimates"],
    ],
    col_widths=[Inches(3.8), Inches(3.0), Inches(5.5)], font_size=14
)

top = SLIDE_H - Inches(1.0)
bullets = [
    (0, "All improvements are additive — the cascade architecture remains the correct foundation"),
]
add_bullet_block(slide, top, bullets, font_size=16)


# ── 21. Closing slide ───────────────────────────────────────────────────────
slide = add_slide()
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, INTEL_BLUE)

add_textbox(slide, MARGIN, Inches(1.5), CONTENT_W, Inches(0.8),
            "Questions?",
            font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide, MARGIN, Inches(2.6), CONTENT_W, Inches(0.5),
            "DINOv2 ViT-Small/14  ·  Two-Stage Cascade  ·  TTA-Averaged Prototypes",
            font_size=18, italic=True, color=RGBColor(0xCC, 0xE5, 0xFF),
            align=PP_ALIGN.CENTER)

for i, (label, value) in enumerate([
    ("Overall Accuracy",   "87.7%"),
    ("Balanced Accuracy",  "0.909"),
    ("Avg. Defect Recall", "91.25%"),
]):
    y = Inches(3.4) + Inches(0.55) * i
    add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.45),
                f"{label}:  {value}",
                font_size=22, bold=(i == 2), color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide, MARGIN, Inches(5.6), CONTENT_W, Inches(0.4),
            "python evaluate.py --cascade --dinov2",
            font_size=14, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────────────────────
out = os.path.join(OUTPUT_DIR, "slides.pptx")
prs.save(out)
print(f"Saved: {out}  ({os.path.getsize(out) // 1024} KB)")
